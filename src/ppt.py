from pptx import Presentation

# https://python-pptx.readthedocs.io/en/latest/user/quickstart.html
prs = Presentation('test.pptx')
text_runs = []
for slide in prs.slides:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                text_runs.append(run.text)

print(text_runs)